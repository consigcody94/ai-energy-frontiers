"""
Quick content QA for the generated deck.
Lists every slide's text frames so we can spot overflow / wrong order / typos
without needing LibreOffice or PowerPoint.

Run with:  python extract_text.py
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = Path(__file__).resolve().parent / "ai_energy_frontiers.pptx"

prs = Presentation(DECK)
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
print()

for idx, slide in enumerate(prs.slides, 1):
    print(f"## Slide {idx}")
    for shape in slide.shapes:
        kind = shape.shape_type
        if shape.has_text_frame:
            tf = shape.text_frame
            # Get position + size
            try:
                left = shape.left / 914400 if shape.left else 0
                top = shape.top / 914400 if shape.top else 0
                w = shape.width / 914400 if shape.width else 0
                h = shape.height / 914400 if shape.height else 0
                pos = f"  ({left:.2f}, {top:.2f}) {w:.2f}x{h:.2f}"
            except Exception:
                pos = ""
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    print(f"  - {text}{pos if not text else ''}")
        elif shape.shape_type == 13:  # picture
            try:
                left = shape.left / 914400
                top = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
                print(f"  [IMAGE] ({left:.2f}, {top:.2f}) {w:.2f}x{h:.2f}")
            except Exception:
                print(f"  [IMAGE]")
    print()
